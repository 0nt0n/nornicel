"""Парсинг корпуса: PDF (PyMuPDF), PPTX (python-pptx), DOCX/DOCM (python-docx) -> чанки.

Чанк = кусок текста + метаданные (doc_id, страница/слайд, язык, структурные подсказки
года/географии из пути). Дальше идёт в извлечение.
"""
import os
import re
import zipfile
from dataclasses import dataclass, asdict
from typing import List, Optional

import config

SUPPORTED_EXT = {".pdf", ".pptx", ".docx", ".docm", ".doc", ".rtf", ".txt"}
_YEAR_RE = re.compile(r"(19|20)\d{2}")
# "Tony_Keating_..." / "Aurelien_Louis_..." — латиница + подчёркивания в начале имени файла,
# типичный паттерн для докладов иностранных спикеров в "Материалы конференций".
_FOREIGN_AUTHOR_RE = re.compile(r"^[A-Za-z]+([_ ][A-Za-z.]+)+_")


@dataclass
class Chunk:
    chunk_id: str
    doc_id: str
    source: str
    page: int
    lang: str
    text: str
    doc_type: str = "unknown"          # верхнеуровневая папка: Доклады/Журналы/Статьи/...
    path_year: Optional[int] = None    # год, если он читается из пути/имени файла
    path_geo: Optional[str] = None     # "foreign", если эвристика по имени сработала, иначе None

    def dict(self):
        return asdict(self)


def _detect_lang(text: str) -> str:
    try:
        from langdetect import detect
        return detect(text[:500])
    except Exception:
        return "unknown"


def _year_from_path(rel_path: str) -> Optional[int]:
    """Ищет 4-значный год в компонентах пути, начиная с самого специфичного (ближе к файлу)."""
    parts = rel_path.replace("\\", "/").split("/")
    for part in reversed(parts):
        m = _YEAR_RE.search(part)
        if m:
            y = int(m.group(0))
            if 1950 <= y <= 2030:
                return y
    return None


def _geo_from_path(rel_path: str, doc_type: str) -> Optional[str]:
    """Слабая эвристика: явно иностранное имя файла в 'Материалы конференций'.
    Возвращает None (а не 'RU'), если сигнала нет — пусть решает LLM по содержимому."""
    name = os.path.basename(rel_path)
    if doc_type == "Материалы конференций" and _FOREIGN_AUTHOR_RE.match(name):
        return "foreign"
    return None


def _doc_id(rel_path: str) -> str:
    """doc_id из относительного пути (не только имени файла) — чтобы избежать коллизий
    одинаковых/похожих имён файлов в разных папках."""
    stem = os.path.splitext(rel_path)[0]
    return stem.replace("\\", "/").replace("/", "__")


def _extract_archives(raw_dir: str) -> None:
    """Разворачивает .zip рядом с самим архивом в <имя>__extracted/, один раз (идемпотентно).
    .rar/.001/.002 не трогаем — нет гарантированной системной зависимости для распаковки."""
    for root, _, files in os.walk(raw_dir):
        if root.endswith("__extracted"):
            continue
        for f in files:
            if not f.lower().endswith(".zip"):
                continue
            zpath = os.path.join(root, f)
            out_dir = os.path.join(root, os.path.splitext(f)[0] + "__extracted")
            if os.path.exists(out_dir):
                continue
            try:
                with zipfile.ZipFile(zpath) as zf:
                    zf.extractall(out_dir)
                print(f"[parse] распакован архив: {os.path.relpath(zpath, raw_dir)}")
            except Exception as e:  # noqa: BLE001
                print(f"[parse] не удалось распаковать {f}: {e}")


def _read_text_file(path: str) -> str:
    """Читает текстовый файл с автоопределением кодировки через chardet."""
    with open(path, "rb") as f:
        raw_bytes = f.read()
    try:
        import chardet  # type: ignore
        detected = chardet.detect(raw_bytes)
        encoding = detected.get("encoding") or "utf-8"
    except ImportError:
        encoding = "utf-8"
    try:
        return raw_bytes.decode(encoding)
    except (UnicodeDecodeError, LookupError):
        return raw_bytes.decode("utf-8", errors="replace")


def _read_rtf(path: str) -> str:
    """Читает .rtf файл через striprtf."""
    text = _read_text_file(path)
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
        return rtf_to_text(text)
    except ImportError:
        # Если striprtf не установлен — возвращаем сырой текст (лучше, чем ничего)
        print(f"[parse] striprtf не установлен, .rtf читается как сырой текст: {path}")
        return text


def _convert_doc_via_libreoffice(path: str) -> Optional[str]:
    """Конвертирует .doc в .docx через LibreOffice headless. Возвращает путь к .docx или None."""
    import subprocess
    import tempfile
    out_dir = tempfile.mkdtemp()
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", out_dir, path],
            capture_output=True, timeout=60
        )
        if result.returncode == 0:
            basename = os.path.splitext(os.path.basename(path))[0] + ".docx"
            converted = os.path.join(out_dir, basename)
            if os.path.exists(converted):
                return converted
    except (FileNotFoundError, subprocess.TimeoutExpired) as e:
        print(f"[parse] LibreOffice не доступен для конвертации .doc: {e}")
    return None


def _read_pages(path: str):
    """Возвращает список (page_number, text) для одного файла любого из поддерживаемых форматов."""
    ext = os.path.splitext(path)[1].lower()
    pages = []
    if ext == ".pdf":
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        for i, page in enumerate(doc, 1):
            pages.append((i, page.get_text()))
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, 1):
            parts = [sh.text for sh in slide.shapes if sh.has_text_frame]
            pages.append((i, "\n".join(parts)))
    elif ext in (".docx", ".docm"):
        from docx import Document
        d = Document(path)
        text = "\n".join(p.text for p in d.paragraphs)
        pages.append((1, text))
    elif ext == ".doc":
        # Сначала пробуем python-docx (иногда .doc — это на самом деле .docx)
        try:
            from docx import Document
            d = Document(path)
            text = "\n".join(p.text for p in d.paragraphs)
            if text.strip():
                pages.append((1, text))
                return pages
        except Exception:
            pass
        # Fallback: конвертация через LibreOffice
        converted = _convert_doc_via_libreoffice(path)
        if converted:
            try:
                from docx import Document
                d = Document(converted)
                text = "\n".join(p.text for p in d.paragraphs)
                pages.append((1, text))
            finally:
                try:
                    os.remove(converted)
                except OSError:
                    pass
        else:
            print(f"[parse] не удалось прочитать .doc: {path}")
    elif ext == ".rtf":
        text = _read_rtf(path)
        pages.append((1, text))
    elif ext == ".txt":
        text = _read_text_file(path)
        pages.append((1, text))
    else:
        raise ValueError(f"Неподдерживаемый формат: {ext}")
    return pages


def _chunk_text(text: str, size: int, overlap: int) -> List[str]:
    text = " ".join(text.split())
    if len(text) <= size:
        return [text] if text.strip() else []
    out, start = [], 0
    while start < len(text):
        end = start + size
        out.append(text[start:end])
        start = end - overlap
    return out


def parse_file(path: str, doc_id: Optional[str] = None, rel_path: Optional[str] = None) -> List[Chunk]:
    rel_path = rel_path or os.path.basename(path)
    doc_id = doc_id or _doc_id(rel_path)
    parts = rel_path.replace("\\", "/").split("/")
    doc_type = parts[0] if len(parts) > 1 else "unknown"
    path_year = _year_from_path(rel_path)
    path_geo = _geo_from_path(rel_path, doc_type)

    chunks: List[Chunk] = []
    for page_no, page_text in _read_pages(path):
        for j, piece in enumerate(_chunk_text(page_text, config.CHUNK_SIZE, config.CHUNK_OVERLAP)):
            chunks.append(Chunk(
                chunk_id=f"{doc_id}::p{page_no}::c{j}",
                doc_id=doc_id,
                source=path,
                page=page_no,
                lang=_detect_lang(piece),
                text=piece,
                doc_type=doc_type,
                path_year=path_year,
                path_geo=path_geo,
            ))
    return chunks


def parse_dir(raw_dir: Optional[str] = None, subdir: Optional[str] = None, skip_processed: bool = True) -> List[Chunk]:
    """Парсит корпус.

    subdir         — ограничиться подпапкой внутри raw_dir (например "Журналы/Цветные металлы/2020").
    skip_processed — не парсить файлы, для которых уже есть завершённый чекпоинт в data/processed/
                      (быстрый путь для "распарсить только новые документы").
    """
    from typing import Dict
    base_raw_dir = raw_dir or config.RAW_DIR
    scan_dir = os.path.join(base_raw_dir, subdir) if subdir else base_raw_dir

    _extract_archives(scan_dir)

    chunks = []
    skipped_ext: Dict[str, int] = {}
    skipped_processed = 0
    for root, _, files in os.walk(scan_dir):
        for f in files:
            ext = os.path.splitext(f)[1].lower()
            path = os.path.join(root, f)
            rel_path = os.path.relpath(path, base_raw_dir)

            if ext not in SUPPORTED_EXT:
                skipped_ext[ext] = skipped_ext.get(ext, 0) + 1
                continue

            doc_id = _doc_id(rel_path)
            if skip_processed:
                ckpt = os.path.join(config.PROCESSED_DIR, f"{doc_id}.json")
                if os.path.exists(ckpt):
                    skipped_processed += 1
                    continue

            try:
                chunks.extend(parse_file(path, doc_id=doc_id, rel_path=rel_path))
            except Exception as e:  # noqa: BLE001
                print(f"[parse] пропущен {rel_path}: {e}")

    if skipped_ext:
        summary = ", ".join(f"{ext or '(без расширения)'}×{n}"
                             for ext, n in sorted(skipped_ext.items(), key=lambda x: -x[1]))
        print(f"[parse] неподдерживаемые типы файлов пропущены: {summary}")
    if skipped_processed:
        print(f"[parse] уже обработано ранее, пропущено файлов: {skipped_processed}")
    return chunks
