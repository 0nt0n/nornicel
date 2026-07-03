"""Парсинг корпуса: PDF (PyMuPDF), PPTX (python-pptx), DOCX (python-docx) -> чанки.

Чанк = кусок текста + метаданные (doc_id, страница/слайд, язык). Дальше идёт в извлечение.
"""
import os
from dataclasses import dataclass, asdict
from typing import List

import config


@dataclass
class Chunk:
    chunk_id: str
    doc_id: str
    source: str
    page: int
    lang: str
    text: str

    def dict(self):
        return asdict(self)


def _detect_lang(text: str) -> str:
    try:
        from langdetect import detect
        return detect(text[:500])
    except Exception:
        return "unknown"


def _chunk_text(text: str, size: int, overlap: int) -> List[str]:
    text = " ".join(text.split())
    if len(text) <= size:
        return [text] if text.strip() else []
    out, start = [], 0
    while start < len(text):
        end = start + size
        out.append(text[start:end])
        start = end - overlap
    return out


def _read_pages(path: str):
    """Возвращает список (page_number, text) для одного файла любого из 3 форматов."""
    ext = os.path.splitext(path)[1].lower()
    pages = []
    if ext == ".pdf":
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        for i, page in enumerate(doc, 1):
            pages.append((i, page.get_text()))
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, 1):
            parts = [sh.text for sh in slide.shapes if sh.has_text_frame]
            pages.append((i, "\n".join(parts)))
    elif ext in (".docx", ".doc"):
        from docx import Document
        d = Document(path)
        text = "\n".join(p.text for p in d.paragraphs)
        pages.append((1, text))
    else:
        raise ValueError(f"Неподдерживаемый формат: {ext}")
    return pages


def parse_file(path: str) -> List[Chunk]:
    doc_id = os.path.splitext(os.path.basename(path))[0]
    chunks: List[Chunk] = []
    for page_no, page_text in _read_pages(path):
        for j, piece in enumerate(_chunk_text(page_text, config.CHUNK_SIZE, config.CHUNK_OVERLAP)):
            chunks.append(Chunk(
                chunk_id=f"{doc_id}::p{page_no}::c{j}",
                doc_id=doc_id,
                source=path,
                page=page_no,
                lang=_detect_lang(piece),
                text=piece,
            ))
    return chunks


def parse_dir(raw_dir: str = None) -> List[Chunk]:
    raw_dir = raw_dir or config.RAW_DIR
    chunks = []
    for root, _, files in os.walk(raw_dir):
        for f in files:
            if os.path.splitext(f)[1].lower() in (".pdf", ".pptx", ".docx", ".doc"):
                try:
                    chunks.extend(parse_file(os.path.join(root, f)))
                except Exception as e:  # noqa: BLE001
                    print(f"[parse] пропущен {f}: {e}")
    return chunks
